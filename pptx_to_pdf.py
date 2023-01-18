import os
from pptx import Presentation
from pdfkit import from_file

# Open the PowerPoint file in path
def pptx_to_pdf(path):
    '''
    this code converts pptx to pdf using the python pptx module

    Then it saves each slide as a separate PDF file using the save method of the slide. After that,
    it uses the pdfkit library to merge all the slide pdfs into one pdf file. Finally, it removes the temporary slide pdfs.
    '''
    prs = Presentation(path)
    # Save each slide as a separate PDF file
    for i in range(prs.slides):
        slide.save('path/to/slide-{}.pdf'.format(i))

    pdf_files = ['path/to/slide-{}.pdf'.format(i) for i in range(len(prs.slides))]
    from_file(pdf_files, 'path/to/pptx.pdf')

    # deleting the residules from partial pdfs
    for pdf_file in pdf_files:
        os.remove(pdf_file)
