from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from pathlib import Path
import PyPDF2
import boto3
#Leo Here, I am thinking of hosting this script onto ec2 that way, once the ec2 gained
#access to our s3 bucket it can just perform the conversion script on the cloud.
#dowload_from_s3 perform the download from s3 bucket with your local credentials
#pdf_to_pptx performs the conversion from pdf to pptx using python pptx module


#download file from s3 with name resource_name, bucket bucket_name, directory path
def dowload_from_s3(resource_name,bucket_name,path):
    '''This is a function that will need to be rewritten once it is hosted on EC2'''
    s3 = boto3.resource(resource_name)
    # Set the name of the S3 bucket and the file you want to download
    bn = bucket_name
    file_name = 'path/to/file.pdf'
    s3.Bucket(bucket_name).download_file(file_name, '/path/to/local/file.pdf')


# Create a PowerPoint presentation with the contents of the pdf
# pptx name = name, save location of the target pdf = file_dir, option = method of conversion

def pdf_to_pptx(name,file_dir,option = 0):
    '''
    There are two ways to perform this conversion, you can
    option 0 pdf to img to pptx conversion
    option 1 pdf to pptx text to pptx conversion
    Right now option 0 is done and tested on my local windows machine with a local pptx
    However option 1 is still work in progress'''
    prs = Presentation()
    if(option == 0 ):
        images = convert_from_path(file_dir)
        # Add the images to the PowerPoint as slides
        for image in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(image, Inches(0), Inches(0), height=Inches(6), width=Inches(8))

        # Save the PowerPoint
        #get the parent directory
        new_path = Path(file_dir).parent.absolute()
        #save the file as name.pptx in the directory
        new_path = new_path.as_posix() + "/" + name + ".pptx"
        #save to new path
        prs.save(new_path)

    # Creating a pdf file object
    if(option > 0):
        #open pdf file directory
        pdfFileObj = open(file_dir,'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        # Getting number of pages in pdf file
        pages = pdfReader.numPages
        # Loop for reading all the Pages
        for i in range(pages):
                pageObj = pdfReader.getPage(i)
                print("Page No: ",i)
                # Extracting text from page
                # And splitting it into chunks of lines
                text = pageObj.extractText().split("  ")
                # Finally the lines are stored into list
                # For iterating over list a loop is used
                for i in range(len(text)):
                        # Printing the line
                        # Lines are seprated using "\n"
                        print(text[i],end="\n\n")
                # For Seprating the Pages
                print()
        # closing the pdf file object
        pdfFileObj.close()
